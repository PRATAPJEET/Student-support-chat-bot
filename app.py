import os
import io
import requests
import streamlit as st
from chatbot import init_chatbot, get_ai_stream_response
from streamlit_lottie import st_lottie
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import PyPDF2

# 1. Page Configuration Setup
st.set_page_config(
    page_title="CampusAI Pro",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize chatbot credentials early and save to client
client = init_chatbot()

# Helper function to fetch asset animations safely
def load_lottie_url(url: str):
    try:
        r = requests.get(url, timeout=5)
        if r.status_code == 200:
            return r.json()
    except:
        return None

# Load interactive floating AI robot animation
lottie_robot = load_lottie_url("https://raw.githubusercontent.com/tuhinsamanta/Lottie-Animation-JSON-URL/main/robot.json")

# Initialize page state
if "current_page" not in st.session_state:
    st.session_state.current_page = "Dashboard Home"

# --- PPTX GENERATION ENGINE ---
def generate_pptx_stream():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_DARK = RGBColor(11, 15, 25)
    LAVENDER = RGBColor(183, 148, 244)
    SOFT_PURPLE = RGBColor(214, 188, 250)
    TEXT_LIGHT = RGBColor(248, 250, 252)
    
    def apply_dark_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    blank_layout = prs.slide_layouts[6]

    # SLIDE 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide1)
    tx_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
    tf = tx_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "🎓 CAMPUSAI PRO"
    p1.font.name, p1.font.size, p1.font.bold, p1.font.color.rgb = 'Helvetica', Pt(54), True, LAVENDER
    p2 = tf.add_paragraph()
    p2.text = "Advanced AI-Powered Multi-Module Student Assistant Terminal"
    p2.font.name, p2.font.size, p2.font.color.rgb = 'Helvetica', Pt(22), SOFT_PURPLE
    p2.space_after = Pt(45)
    p3 = tf.add_paragraph()
    p3.text = "Presented by: Jeet Pratap (B.Tech CSE Enthusiast)\nProject Mentorship: Nikhil Sharma"
    p3.font.name, p3.font.size, p3.font.color.rgb = 'Helvetica', Pt(16), TEXT_LIGHT

    # Save to stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# 2. Sidebar UI Navigation
st.sidebar.title("🎓 CampusAI")
st.sidebar.caption("Your Intelligent Student Assistant")
st.sidebar.markdown("---")

if st.session_state.current_page != "Dashboard Home":
    if st.sidebar.button("🏠 Back to Home Dashboard", use_container_width=True, type="primary"):
        st.session_state.current_page = "Dashboard Home"
        st.rerun()
    st.sidebar.markdown("---")

st.sidebar.subheader("🛠️ Navigation")
pages = [
    "Dashboard Home", 
    "AI Chat Thread", 
    "Student Success Support", 
    "Campus Knowledge Base", 
    "CS Programming Help", 
    "Chat with PDF Notes", 
    "Voice Audio Assistant"
]

for page in pages:
    label = f"💜 {page}" if st.session_state.current_page == page else page
    if st.sidebar.button(label, use_container_width=True):
        st.session_state.current_page = page

# 3. Premium Cyber-Dark Lavender Style Injection
st.markdown("""
    <style>
        .stApp, [data-testid="stAppViewContainer"] {
            background-color: #0B0F19 !important;
            color: #F8FAFC !important;
        }
        [data-testid="stSidebar"] { background-color: #0F172A !important; border-right: 1px solid #1E293B !important; }
        .hero-title { font-size: 3.2rem !important; font-weight: 800 !important; color: #B794F4 !important; margin-bottom: 5px !important; }
        .hero-subtitle { font-size: 2rem !important; font-weight: 700 !important; color: #F8FAFC !important; margin-bottom: 20px !important; }
        .hero-desc { font-size: 1.15rem !important; color: #94A3B8 !important; line-height: 1.7 !important; margin-bottom: 30px !important; }
        
        .about-card { 
            background: #1E293B !important; 
            border-left: 5px solid #B794F4 !important; 
            padding: 20px !important; 
            border-radius: 12px !important; 
            margin-bottom: 25px !important; 
        }
        
        .feature-card { background: #1E293B !important; padding: 30px 25px !important; border-radius: 20px !important; text-align: center !important; border: 1px solid #334155 !important; margin-bottom: 15px !important; transition: all 0.3s; }
        .feature-card:hover { transform: translateY(-5px); border-color: #B794F4 !important; }
        .feature-icon { font-size: 2.5rem !important; margin-bottom: 15px !important; }
        .feature-title { font-size: 1.35rem !important; font-weight: 700 !important; color: #F8FAFC !important; margin-bottom: 10px !important; }
        .feature-desc { font-size: 0.95rem !important; color: #94A3B8 !important; line-height: 1.5 !important; }
    </style>
""", unsafe_allow_html=True)

# 4. Page Routing Logic
if st.session_state.current_page == "Dashboard Home":
    col1, col2 = st.columns([1.3, 1], vertical_alignment="center")
    with col1:
        st.markdown('<div class="hero-title">AI Student Assistant</div>', unsafe_allow_html=True)
        st.markdown('<div class="hero-subtitle">Your Personal AI Learning Companion</div>', unsafe_allow_html=True)
        st.markdown('<div class="hero-desc">Learn faster with an intelligent AI assistant framework that can answer your academic questions, summarize notes, generate quizzes, and prepare smart study plans.</div>', unsafe_allow_html=True)
        
        # About Me Card
        st.markdown("""
        <div class="about-card">
            <span style="font-weight:700; font-size:1.1rem; display:block; margin-bottom:5px; color: #D6BCFA;">💡 About Me</span>
            <span style="font-size:0.98rem; opacity:0.9;">I am <strong>Jeet Pratap</strong>, a B.Tech CSE enthusiast passionate about engineering custom AI pipelines, context indexing engines, and modern UI dashboard frameworks.</span>
        </div>
        """, unsafe_allow_html=True)
        
        btn_c1, btn_c2 = st.columns([1, 1])
        with btn_c1:
            if st.button("🔮 Start Chat Now", use_container_width=True):
                st.session_state.current_page = "AI Chat Thread"
                st.rerun()
        with btn_c2:
            with st.spinner("Preparing PowerPoint file..."):
                ppt_data = generate_pptx_stream()
            st.download_button(label="📺 Download Presentation.pptx", data=ppt_data, file_name="CampusAI_Presentation.pptx", use_container_width=True)

    with col2:
        if lottie_robot:
            st_lottie(lottie_robot, height=380, speed=1.2, loop=True, key="premium_robot")
            
    st.markdown("<br><br><h2 style='text-align: center; color: #D6BCFA;'>Modules</h2>", unsafe_allow_html=True)
    
    # Grid Row 1 (Modules 1-3)
    grid1, grid2, grid3 = st.columns(3)
    
    with grid1:
        st.markdown('<div class="feature-card"><div class="feature-icon">🔮</div><div class="feature-title">AI Chatbot</div></div>', unsafe_allow_html=True)
        if st.button("Open AI Chat Thread", key="go_chat", use_container_width=True): 
            st.session_state.current_page = "AI Chat Thread"
            st.rerun()
            
    with grid2:
        st.markdown('<div class="feature-card"><div class="feature-icon">🤝</div><div class="feature-title">Success Advisor</div></div>', unsafe_allow_html=True)
        if st.button("Open Success Pipeline", key="go_success", use_container_width=True): 
            st.session_state.current_page = "Student Success Support"
            st.rerun()
            
    with grid3:
        st.markdown('<div class="feature-card"><div class="feature-icon">🏢</div><div class="feature-title">Knowledge System</div></div>', unsafe_allow_html=True)
        if st.button("Open Institutional Base", key="go_kb", use_container_width=True): 
            st.session_state.current_page = "Campus Knowledge Base"
            st.rerun()

    st.markdown("<br>", unsafe_allow_html=True)
    
    # Grid Row 2 (Modules 4-6)
    grid4, grid5, grid6 = st.columns(3)
    
    with grid4:
        st.markdown('<div class="feature-card"><div class="feature-icon">💻</div><div class="feature-title">CS Programming TA</div></div>', unsafe_allow_html=True)
        if st.button("Open Coding TA Engine", key="go_cs", use_container_width=True): 
            st.session_state.current_page = "CS Programming Help"
            st.rerun()
            
    with grid5:
        st.markdown('<div class="feature-card"><div class="feature-icon">📄</div><div class="feature-title">Notes Summarizer</div></div>', unsafe_allow_html=True)
        if st.button("Open Document Processor", key="go_pdf", use_container_width=True): 
            st.session_state.current_page = "Chat with PDF Notes"
            st.rerun()
            
    with grid6:
        st.markdown('<div class="feature-card"><div class="feature-icon">🎙️</div><div class="feature-title">Audio Assistant</div></div>', unsafe_allow_html=True)
        if st.button("Open Voice Simulator", key="go_voice", use_container_width=True): 
            st.session_state.current_page = "Voice Audio Assistant"
            st.rerun()

# --- PAGE 1: AI Chat Thread ---
elif st.session_state.current_page == "AI Chat Thread":
    st.title("💬 General AI Chat Thread")
    if "messages" not in st.session_state: 
        st.session_state.messages = []
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]): 
            st.markdown(msg["content"])
        
    if prompt := st.chat_input("Ask anything..."):
        with st.chat_message("user"): 
            st.markdown(prompt)
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("assistant"):
            placeholder = st.empty()
            resp = ""
            # FIXED: Formatted layout context and passed all 3 required arguments
            formatted_history = [{"role": m["role"], "parts": [m["content"]]} for m in st.session_state.messages[:-1]]
            for chunk in get_ai_stream_response(client, formatted_history, prompt):
                resp += chunk
                placeholder.markdown(resp)
        st.session_state.messages.append({"role": "assistant", "content": resp})

# --- PAGE 2: Student Success Support ---
elif st.session_state.current_page == "Student Success Support":
    st.title("🤝 Student Success Advisor")
    st.info("I am optimized to help you with college pipelines, deadlines, and tracking resources.")
    if "success_msgs" not in st.session_state: 
        st.session_state.success_msgs = []
    
    for msg in st.session_state.success_msgs:
        with st.chat_message(msg["role"]): 
            st.markdown(msg["content"])
        
    if prompt := st.chat_input("Ask about university guidelines or deadlines..."):
        with st.chat_message("user"): 
            st.markdown(prompt)
        st.session_state.success_msgs.append({"role": "user", "content": prompt})
        
        context_prompt = f"As a highly knowledgeable University Student Success Advisor, answer this: {prompt}"
        with st.chat_message("assistant"):
            placeholder = st.empty()
            resp = ""
            formatted_history = [{"role": m["role"], "parts": [m["content"]]} for m in st.session_state.success_msgs[:-1]]
            for chunk in get_ai_stream_response(client, formatted_history, context_prompt):
                resp += chunk
                placeholder.markdown(resp)
        st.session_state.success_msgs.append({"role": "assistant", "content": resp})

# --- PAGE 3: CS Programming Help ---
elif st.session_state.current_page == "CS Programming Help":
    st.title("💻 CS Programming Assistant TA")
    st.info("Paste your broken code or logic questions here. I will help you debug!")
    if "cs_msgs" not in st.session_state: 
        st.session_state.cs_msgs = []
    
    for msg in st.session_state.cs_msgs:
        with st.chat_message(msg["role"]): 
            st.markdown(msg["content"])
        
    if prompt := st.chat_input("Paste code here..."):
        with st.chat_message("user"): 
            st.markdown(prompt)
        st.session_state.cs_msgs.append({"role": "user", "content": prompt})
        
        context_prompt = f"Act as an expert Computer Science Teaching Assistant. Evaluate this query, explain bugs, and provide corrected code: {prompt}"
        with st.chat_message("assistant"):
            placeholder = st.empty()
            resp = ""
            formatted_history = [{"role": m["role"], "parts": [m["content"]]} for m in st.session_state.cs_msgs[:-1]]
            for chunk in get_ai_stream_response(client, formatted_history, context_prompt):
                resp += chunk
                placeholder.markdown(resp)
        st.session_state.cs_msgs.append({"role": "assistant", "content": resp})

# --- PAGE 4: Campus Knowledge Base ---
elif st.session_state.current_page == "Campus Knowledge Base":
    st.title("🏢 Campus Knowledge Base")
    st.info("Query standard institutional rules, library hours, and compliance policies.")
    if "kb_msgs" not in st.session_state: 
        st.session_state.kb_msgs = []
    
    for msg in st.session_state.kb_msgs:
        with st.chat_message(msg["role"]): 
            st.markdown(msg["content"])
        
    if prompt := st.chat_input("Search the campus index..."):
        with st.chat_message("user"): 
            st.markdown(prompt)
        st.session_state.kb_msgs.append({"role": "user", "content": prompt})
        
        context_prompt = f"Act as a formal Campus Knowledge Base Indexer. Provide clear, policy-based answers: {prompt}"
        with st.chat_message("assistant"):
            placeholder = st.empty()
            resp = ""
            formatted_history = [{"role": m["role"], "parts": [m["content"]]} for m in st.session_state.kb_msgs[:-1]]
            for chunk in get_ai_stream_response(client, formatted_history, context_prompt):
                resp += chunk
                placeholder.markdown(resp)
        st.session_state.kb_msgs.append({"role": "assistant", "content": resp})

# --- PAGE 5: Chat with PDF Notes ---
elif st.session_state.current_page == "Chat with PDF Notes":
    st.title("📄 Notes Summarizer & PDF Processor")
    st.markdown("Upload your syllabus or lecture notes and chat directly with your document.")
    
    uploaded_file = st.file_uploader("Upload a PDF document", type="pdf")
    
    if uploaded_file is not None:
        if "pdf_text" not in st.session_state or st.session_state.get("pdf_name") != uploaded_file.name:
            with st.spinner("Extracting text from PDF..."):
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                st.session_state.pdf_text = text
                st.session_state.pdf_name = uploaded_file.name
                st.success("✅ Document processed successfully!")
                
        if "pdf_msgs" not in st.session_state: 
            st.session_state.pdf_msgs = []
        for msg in st.session_state.pdf_msgs:
            with st.chat_message(msg["role"]): 
                st.markdown(msg["content"])
            
        if prompt := st.chat_input("Ask a question about the uploaded document..."):
            with st.chat_message("user"): 
                st.markdown(prompt)
            st.session_state.pdf_msgs.append({"role": "user", "content": prompt})
            
            context_prompt = f"Based strictly on the following document text, answer the user's question.\n\nDocument Text:\n{st.session_state.pdf_text[:5000]}...\n\nUser Question: {prompt}"
            
            with st.chat_message("assistant"):
                placeholder = st.empty()
                resp = ""
                formatted_history = [{"role": m["role"], "parts": [m["content"]]} for m in st.session_state.pdf_msgs[:-1]]
                for chunk in get_ai_stream_response(client, formatted_history, context_prompt):
                    resp += chunk
                    placeholder.markdown(resp)
            st.session_state.pdf_msgs.append({"role": "assistant", "content": resp})

# --- PAGE 6: Voice Audio Assistant ---
elif st.session_state.current_page == "Voice Audio Assistant":
    st.title("🎙️ Voice Audio Assistant")
    st.info("Record a voice query right from your browser.")
    
    audio_value = st.audio_input("Record your question here")
    
    if audio_value:
        st.success("✅ Audio captured successfully! You can review your audio below.")
        st.audio(audio_value)